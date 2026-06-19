"""Generate the ABI Ten Websites pptx deck."""
import json, os, urllib.request
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
CDN = "https://abi-assets.vercel.app"
SITES = json.loads((ROOT / "_content" / "sites.json").read_text())["sites"]
MANIFEST = json.loads((ROOT / "_deploy" / "manifest.json").read_text())

# Download all 10 logos for embedding
LOGO_DIR = ROOT / "_deploy" / "pptx_logos"
LOGO_DIR.mkdir(exist_ok=True, parents=True)
for s in SITES:
    dst = LOGO_DIR / s["logo"]
    if not dst.exists():
        try:
            urllib.request.urlretrieve(f"{CDN}/logos/{s['logo']}", dst)
        except Exception as e:
            print(f"  WARN: couldn't fetch logo for {s['slug']}: {e}")

def hexrgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
PURE_BG = hexrgb("#0a0c12")
INK = hexrgb("#f1f3ff")
MUT = hexrgb("#8b91b8")

def add_solid_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Helvetica"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb

# ---------- Slide 1: Title ----------
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(s1, PURE_BG)
# accent stripe
stripe = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
stripe.fill.solid(); stripe.fill.fore_color.rgb = hexrgb("#00f0ff")
stripe.line.fill.background()
add_text(s1, Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
         "AMERICAN BARBER INSTITUTE  ·  TEN VISIONS", size=11, bold=True, color=hexrgb("#00f0ff"))
add_text(s1, Inches(0.6), Inches(1.4), Inches(12), Inches(2.2),
         "One brand.\nTen distinct websites.", size=64, bold=True, color=INK)
add_text(s1, Inches(0.6), Inches(4.4), Inches(11), Inches(1.5),
         "Ten production-ready website concepts for ABI — NYC's only dedicated barber school.\nEach with its own visual identity. Built on the same real content. Deployed independently.",
         size=18, color=MUT)
add_text(s1, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
         "Mobile-first  ·  SEO + JSON-LD  ·  Multi-page  ·  abi-app-1 → abi-app-10.vercel.app",
         size=12, color=hexrgb("#00f0ff"))

# ---------- Slide 2: Overview ----------
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(s2, PURE_BG)
add_text(s2, Inches(0.6), Inches(0.5), Inches(12), Inches(0.5),
         "AT A GLANCE", size=11, bold=True, color=hexrgb("#00f0ff"))
add_text(s2, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
         "Ten unique deployments", size=36, bold=True, color=INK)

# 5x2 grid of mini-cards
card_w = Inches(2.4); card_h = Inches(2.4); gx = Inches(0.6); gy = Inches(2.1); gap = Inches(0.15)
for i, s in enumerate(SITES):
    col = i % 5
    row = i // 5
    x = gx + (card_w + gap) * col
    y = gy + (card_h + gap) * row
    card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.adjustments[0] = 0.08
    card.fill.solid(); card.fill.fore_color.rgb = hexrgb("#13161f")
    card.line.color.rgb = hexrgb(s["primary_color"])
    card.line.width = Pt(0.75)
    # accent strip
    strip = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.12))
    strip.fill.solid(); strip.fill.fore_color.rgb = hexrgb(s["primary_color"])
    strip.line.fill.background()
    # number
    add_text(s2, x + Inches(0.18), y + Inches(0.22), card_w - Inches(0.36), Inches(0.45),
             f"0{i+1}" if i<9 else "10", size=14, bold=True, color=hexrgb(s["primary_color"]))
    # name
    add_text(s2, x + Inches(0.18), y + Inches(0.72), card_w - Inches(0.36), Inches(0.55),
             s["name"], size=18, bold=True, color=INK)
    # theme
    add_text(s2, x + Inches(0.18), y + Inches(1.25), card_w - Inches(0.36), Inches(0.5),
             s["theme_word"], size=10, color=MUT)
    # URL
    add_text(s2, x + Inches(0.18), y + Inches(1.95), card_w - Inches(0.36), Inches(0.4),
             f"{s['vercel_name']}.vercel.app", size=10, bold=True, color=hexrgb(s["primary_color"]))

# ---------- Slides 3-12: per site ----------
for i, s in enumerate(SITES):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    accent = hexrgb(s["primary_color"])
    add_solid_bg(sl, PURE_BG)
    # left accent panel
    panel = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.6), prs.slide_height)
    panel.fill.solid(); panel.fill.fore_color.rgb = hexrgb("#0f121b")
    panel.line.fill.background()
    # top thin accent
    top = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top.fill.solid(); top.fill.fore_color.rgb = accent
    top.line.fill.background()
    # number
    add_text(sl, Inches(0.6), Inches(0.5), Inches(4), Inches(0.5),
             f"CONCEPT 0{i+1}/10" if i<9 else "CONCEPT 10/10", size=11, bold=True, color=accent)
    # logo (large on left panel)
    logo_path = LOGO_DIR / s["logo"]
    if logo_path.exists():
        try:
            sl.shapes.add_picture(str(logo_path), Inches(0.7), Inches(1.4), height=Inches(2.6))
        except Exception:
            pass
    # name (huge)
    add_text(sl, Inches(0.6), Inches(4.5), Inches(4), Inches(1.5),
             s["name"], size=46, bold=True, color=INK)
    # theme
    add_text(sl, Inches(0.6), Inches(5.7), Inches(4), Inches(0.5),
             s["theme_word"], size=14, color=MUT)
    # URL pill
    pill = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.6), Inches(3.6), Inches(0.5))
    pill.adjustments[0] = 0.5
    pill.fill.solid(); pill.fill.fore_color.rgb = accent
    pill.line.fill.background()
    add_text(sl, Inches(0.6), Inches(6.66), Inches(3.6), Inches(0.4),
             f"{s['vercel_name']}.vercel.app", size=12, bold=True, color=PURE_BG, align=PP_ALIGN.CENTER)

    # right side: highlights
    add_text(sl, Inches(5.2), Inches(0.9), Inches(7.5), Inches(0.5),
             "Concept signature", size=11, bold=True, color=accent)
    add_text(sl, Inches(5.2), Inches(1.3), Inches(7.5), Inches(1.3),
             s["description"], size=22, bold=True, color=INK)
    add_text(sl, Inches(5.2), Inches(2.9), Inches(7.5), Inches(0.4),
             "Tagline variant", size=11, bold=True, color=accent)
    add_text(sl, Inches(5.2), Inches(3.3), Inches(7.5), Inches(0.7),
             f'"{s["tagline_variant"]}"', size=18, color=INK)
    add_text(sl, Inches(5.2), Inches(4.2), Inches(7.5), Inches(0.4),
             "Pages", size=11, bold=True, color=accent)
    add_text(sl, Inches(5.2), Inches(4.55), Inches(7.5), Inches(2.2),
             "•  Home — hero, programs, careers, reviews, tuition, locations, FAQ\n•  Programs — every program, schedule, tuition plan\n•  Gallery — 30+ photos from the clinic floor\n•  FAQ — 15 questions with full answers (JSON-LD)\n•  Contact — both campuses + lead form",
             size=13, color=MUT)
    add_text(sl, Inches(5.2), Inches(6.6), Inches(7.5), Inches(0.5),
             "Mobile-first  •  Hamburger nav  •  SEO meta + OpenGraph  •  Sitemap + robots  •  CDN-backed assets",
             size=11, color=MUT)

# ---------- Final slide: All 10 URLs ----------
sf = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sf, PURE_BG)
stripe2 = sf.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
stripe2.fill.solid(); stripe2.fill.fore_color.rgb = hexrgb("#d4af37"); stripe2.line.fill.background()
add_text(sf, Inches(0.6), Inches(0.5), Inches(12), Inches(0.5),
         "ALL LIVE URLS", size=11, bold=True, color=hexrgb("#d4af37"))
add_text(sf, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
         "Ten sites. One naming pattern.", size=38, bold=True, color=INK)
add_text(sf, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
         "abi-app-N.vercel.app — change N from 1 to 10.", size=14, color=MUT)

col1_x = Inches(0.6); col2_x = Inches(7.0); list_y = Inches(2.8)
for i, s in enumerate(SITES):
    x = col1_x if i < 5 else col2_x
    y = list_y + Inches(0.7 * (i % 5))
    # number bullet
    add_text(sf, x, y, Inches(0.5), Inches(0.6),
             f"{i+1:02d}.", size=22, bold=True, color=hexrgb(s["primary_color"]))
    add_text(sf, x + Inches(0.55), y, Inches(2.0), Inches(0.4),
             s["name"], size=15, bold=True, color=INK)
    add_text(sf, x + Inches(0.55), y + Inches(0.32), Inches(5.4), Inches(0.4),
             f"https://{s['vercel_name']}.vercel.app",
             size=12, color=hexrgb(s["primary_color"]))

add_text(sf, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
         "Repo: github.com/kazi-reprime/ABI-10-Websites   ·   Assets CDN: abi-assets.vercel.app",
         size=12, color=MUT)

out = ROOT / "_deploy" / "ABI-Ten-Websites.pptx"
prs.save(out)
print(f"PPTX saved: {out}")
print(f"Size: {out.stat().st_size//1024} KB")
print(f"Slides: {len(prs.slides)}")
